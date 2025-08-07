import os
import pandas as pd
from docx import Document
import pdfplumber
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
import streamlit as st
import faiss
import numpy as np
import cohere
from openai import OpenAI
from dotenv import load_dotenv
import sqlite3
import qrcode
from io import BytesIO
from datetime import datetime, timedelta, date, timezone
import pickle
import json
import re
from netmiko import ConnectHandler
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import TextLoader
from langchain_community.vectorstores import FAISS as LangchainFAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_openai import ChatOpenAI
from langchain.chains import RetrievalQA, LLMChain  # LLMChain 추가
from langchain.prompts import PromptTemplate
from google.oauth2 import service_account
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from langchain.globals import set_llm_cache
from langchain.cache import SQLiteCache

# --- 페이지 기본 설정 ---
st.set_page_config(page_title="유지보수 챗봇", layout="wide")

# LangChain의 LLM 호출 결과를 캐싱하여 비용과 속도를 최적화
set_llm_cache(SQLiteCache(database_path="langchain_llm_cache.sqlite"))

# .env 파일 로드
load_dotenv()

# --- 상수 및 설정 ---
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
COHERE_API_KEY = os.getenv('COHERE_API_KEY')

if not OPENAI_API_KEY or not COHERE_API_KEY:
    st.error("OPENAI_API_KEY 또는 COHERE_API_KEY가 .env 파일에 설정되지 않았습니다.")
    st.stop()

# 폴더 및 파일 경로
UPLOAD_FOLDER = 'Uploads'
ASSET_DB_PATH = 'assets.db'
GOOGLE_CREDS_PATH = 'credentials.json'
GOOGLE_TOKEN_PATH = 'token.pickle'
FAISS_INDEX_PATH = 'faiss_index.bin'
FAISS_METADATA_PATH = 'faiss_metadata.pkl'
NETWORK_RAW_DATA_FILE = "network_raw_data.txt"
NETWORK_FAISS_INDEX_PATH = "network_faiss_index"

# 허용된 파일 확장자 및 최대 파일 크기 (16MB)
ALLOWED_EXTENSIONS = {'csv', 'docx', 'pdf', 'jpg', 'jpeg', 'png', 'ppt', 'pptx', 'txt'}
MAX_FILE_SIZE = 16 * 1024 * 1024

# 네트워크 챗봇 프롬프트 템플릿
NETWORK_PROMPT_TEMPLATE = """
### 역할 및 목표
당신은 숙련된 네트워크 엔지니어의 역할을 수행하는 AI 어시스턴트입니다. 당신의 목표는 제공된 '컨텍스트'(Cisco 장비의 다양한 show 명령어 결과)를 깊이 있게 분석하고, 사용자의 질문에 대해 명확하고 구조화된 답변을 제공하는 것입니다.

### 행동 지침
1. **정보 종합:** 사용자의 질문에 대한 답은 여러 명령어 결과에 흩어져 있을 수 있습니다. 예를 들어, 특정 VLAN의 IP 정보를 알려달라는 요청에는 `show vlan`, `show ip interface brief`, `show running-config`의 내용을 모두 종합하여 답변해야 합니다.
2. **정확한 정보 추출:** `show running-config`의 `interface VlanX` 섹션에서 `ip address [IP 주소] [서브넷 마스크]` 형식의 정보를 정확히 찾아내세요.
3. **구조화된 답변:** 가능하다면, 정보를 표(Markdown 테이블 형식)로 정리하여 가독성을 높여주세요.
4. **추론 및 요약:** 단순히 텍스트를 복사하지 말고, 질문의 의도에 맞게 정보를 요약하고 재구성하여 답변하세요.
5. **제한된 정보 내에서 답변:** 주어진 '컨텍스트'에 정보가 없는 경우, "제공된 로그 정보에서는 해당 내용을 찾을 수 없습니다."라고 명확하게 답변하세요.

### 컨텍스트
{context}

### 질문
{question}

### 답변 (위 지침에 따라 생성):
"""
NETWORK_PROMPT = PromptTemplate(
    template=NETWORK_PROMPT_TEMPLATE, input_variables=["context", "question"]
)

# --- 초기화 ---
@st.cache_resource
def get_clients():
    co_client = cohere.Client(COHERE_API_KEY)
    openai_client = OpenAI(api_key=OPENAI_API_KEY)
    dimension = 1024
    index = faiss.IndexFlatL2(dimension)
    metadata = []
    if os.path.exists(FAISS_INDEX_PATH) and os.path.exists(FAISS_METADATA_PATH):
        try:
            index = faiss.read_index(FAISS_INDEX_PATH)
            with open(FAISS_METADATA_PATH, 'rb') as f:
                metadata = pickle.load(f)
            validated_metadata = []
            for item in metadata:
                if isinstance(item, dict) and 'text' in item:
                    validated_metadata.append(item)
                else:
                    st.warning(f"잘못된 메타데이터 항목 발견: {item}. 무시됩니다.")
            metadata = validated_metadata
            if index.ntotal != len(metadata):
                st.warning("FAISS 인덱스와 메타데이터 길이가 일치하지 않습니다. 새 인덱스를 생성합니다.")
                index = faiss.IndexFlatL2(dimension)
                metadata = []
                if os.path.exists(FAISS_INDEX_PATH):
                    os.remove(FAISS_INDEX_PATH)
                if os.path.exists(FAISS_METADATA_PATH):
                    os.remove(FAISS_METADATA_PATH)
        except Exception as e:
            st.warning(f"FAISS 인덱스 로드 중 오류: {e}. 새 인덱스를 생성합니다.")
            index = faiss.IndexFlatL2(dimension)
            metadata = []
    return co_client, index, openai_client, metadata

@st.cache_resource
def load_network_rag_chain():
    if not os.path.exists(NETWORK_FAISS_INDEX_PATH):
        return None
    try:
        embeddings = HuggingFaceEmbeddings(model_name="BM-K/KoSimCSE-roberta-multitask")
        vectorstore = LangchainFAISS.load_local(NETWORK_FAISS_INDEX_PATH, embeddings, allow_dangerous_deserialization=True)
        retriever = vectorstore.as_retriever()
        llm = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=0, openai_api_key=OPENAI_API_KEY)
        return RetrievalQA.from_chain_type(
            llm=llm, chain_type="stuff", retriever=retriever,
            return_source_documents=True, chain_type_kwargs={"prompt": NETWORK_PROMPT}
        )
    except Exception as e:
        st.error(f"네트워크 RAG 체인 로드 중 오류: {e}")
        return None

co, index, openai_client, metadata = get_clients()

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# --- 네트워크 데이터 수집 및 인덱싱 ---
def fetch_network_data(device_info, commands_to_run):
    all_output = ""
    try:
        with ConnectHandler(**device_info) as net_connect:
            net_connect.enable()
            net_connect.fast_cli = True  # fast_cli 모드 활성화
            for cmd in commands_to_run:
                # send_command_timing을 사용하여 더 빠른 응답을 기대
                output = net_connect.send_command_timing(cmd, delay_factor=1, max_loops=150, read_timeout=60)
                all_output += f"\n--- {cmd} ---\n{output}\n"
            net_connect.fast_cli = False # fast_cli 모드 비활성화
        with open(NETWORK_RAW_DATA_FILE, "w", encoding="utf-8") as f:
            f.write(all_output)
        return True, None
    except Exception as e:
        error_message = f"데이터 수집 중 오류 발생: {e}"
        return False, error_message

def build_network_vector_store():
    if not os.path.exists(NETWORK_RAW_DATA_FILE):
        error_message = f"{NETWORK_RAW_DATA_FILE} 파일이 없습니다. 먼저 데이터를 수집하세요."
        return False, error_message
    try:
        loader = TextLoader(NETWORK_RAW_DATA_FILE, encoding="utf-8")
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=150)
        docs = text_splitter.split_documents(documents)
        embeddings = HuggingFaceEmbeddings(model_name="BM-K/KoSimCSE-roberta-multitask")
        vectorstore = LangchainFAISS.from_documents(docs, embeddings)
        vectorstore.save_local(NETWORK_FAISS_INDEX_PATH)
        return True, None
    except Exception as e:
        error_message = f"네트워크 벡터 스토어 구축 중 오류 발생: {e}"
        return False, error_message

# --- LangChain 날짜 추출 체인 ---
@st.cache_resource
def get_date_extraction_chain():
    llm = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=0, openai_api_key=OPENAI_API_KEY)
    date_extraction_template = """사용자의 질문에서 날짜 또는 기간 정보를 추출하여 JSON 형식으로 반환해주세요.
- 오늘 날짜는 {today} 입니다.
- '오늘', '내일', '이번 주' 등 상대적 표현은 정확한 날짜로 변환해야 합니다.
- 날짜는 'YYYY-MM-DD' 형식이어야 합니다.
- 날짜 정보가 없으면, 오늘 날짜를 start_date와 end_date로 설정합니다.
- 질문에서 검색할 키워드(사람 이름, 이벤트 종류 등)를 'query' 키로 추출합니다. 키워드가 없거나 '전체', '모든' 등의 단어만 있으면 '전체'로 설정합니다.

---

**예시**
오늘 날짜: 2025-08-01

질문: 8월 15일 전체 일정을 알려줘
JSON: {{"start_date": "2025-08-15", "end_date": "2025-08-15", "query": "전체"}}

질문: 다음 주 백종윤의 일정이 어떻게 돼?
JSON: {{"start_date": "2025-08-04", "end_date": "2025-08-10", "query": "백종윤"}}

질문: 현우 미팅
JSON: {{"start_date": "2025-08-01", "end_date": "2025-08-01", "query": "현우 미팅"}}

---

오늘 날짜: {today}
질문: {question}
JSON:"""
    date_extraction_prompt = PromptTemplate(template=date_extraction_template, input_variables=["today", "question"])
    return LLMChain(llm=llm, prompt=date_extraction_prompt)

date_extraction_chain = get_date_extraction_chain()

# --- SQLite 자산관리 기능 ---
def init_asset_db():
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS assets (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            "자산번호" TEXT NOT NULL UNIQUE,
            "위치" TEXT,
            "장비상태" TEXT,
            "제조사" TEXT,
            "시리얼" TEXT,
            "구분" TEXT,
            "모델명" TEXT,
            "기타" TEXT
        )
    ''')
    c.execute("PRAGMA table_info(assets)")
    columns = [info[1] for info in c.fetchall()]
    if "시스템명" in columns and "기타" not in columns:
        c.execute('''
            CREATE TABLE assets_new (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                "자산번호" TEXT NOT NULL UNIQUE,
                "위치" TEXT,
                "장비상태" TEXT,
                "제조사" TEXT,
                "시리얼" TEXT,
                "구분" TEXT,
                "모델명" TEXT,
                "기타" TEXT
            )
        ''')
        c.execute('''
            INSERT INTO assets_new (id, "자산번호", "위치", "장비상태", "제조사", "시리얼", "구분", "모델명", "기타")
            SELECT id, "자산번호", "위치", "장비상태", "제조사", "시리얼", "구분", "모델명", "시스템명"
            FROM assets
        ''')
        c.execute("DROP TABLE assets")
        c.execute("ALTER TABLE assets_new RENAME TO assets")
    conn.commit()
    conn.close()

def add_asset(data):
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    try:
        c.execute('''
            INSERT INTO assets ("자산번호", "위치", "장비상태", "제조사", "시리얼", "구분", "모델명", "기타")
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)''', (data['자산번호'], data['위치'], data['장비상태'], data['제조사'], data['시리얼'], data['구분'], data['모델명'], data['기타']))
        conn.commit()
        return True
    except sqlite3.IntegrityError as e:
        st.error(f"DB 삽입 오류: {e} (자산번호: {data['자산번호']})")
        return False
    except Exception as e:
        st.error(f"DB 삽입 중 예상치 못한 오류: {e}")
        return False
    finally:
        conn.close()

def get_assets():
    conn = sqlite3.connect(ASSET_DB_PATH)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()
    c.execute('SELECT id, "자산번호", "위치", "장비상태", "제조사", "시리얼", "구분", "모델명", "기타" FROM assets ORDER BY id DESC')
    assets = [dict(row) for row in c.fetchall()]
    conn.close()
    return assets

def get_unique_statuses():
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    c.execute('SELECT DISTINCT "장비상태" FROM assets WHERE "장비상태" IS NOT NULL')
    statuses = [row[0] for row in c.fetchall()]
    conn.close()
    default_statuses = ["사용중", "미사용", "수리중", "폐기", "정상", "비정상"]
    unique_statuses = list(set(default_statuses + statuses))
    return unique_statuses

def update_asset(data):
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    try:
        c.execute('''
            UPDATE assets SET
            "자산번호"=?, "위치"=?, "장비상태"=?, "제조사"=?, "시리얼"=?, "구분"=?, "모델명"=?, "기타"=?
            WHERE id=?
        ''', (data['자산번호'], data['위치'], data['장비상태'], data['제조사'], data['시리얼'], data['구분'], data['모델명'], data['기타'], data['id']))
        conn.commit()
    except sqlite3.IntegrityError as e:
        st.error(f"DB 수정 오류: {e} (자산번호: {data['자산번호']})")
    except Exception as e:
        st.error(f"DB 수정 중 예상치 못한 오류: {e}")
    finally:
        conn.close()

def delete_asset(asset_id):
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    try:
        c.execute("DELETE FROM assets WHERE id=?", (asset_id,))
        conn.commit()
    except Exception as e:
        st.error(f"DB 삭제 오류: {e}")
    finally:
        conn.close()

def delete_all_assets():
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    try:
        c.execute("DELETE FROM assets")
        c.execute("DELETE FROM sqlite_sequence WHERE name='assets'")
        conn.commit()
    except Exception as e:
        st.error(f"DB 초기화 오류: {e}")
    finally:
        conn.close()

def get_existing_asset_numbers():
    conn = sqlite3.connect(ASSET_DB_PATH)
    c = conn.cursor()
    c.execute('SELECT "자산번호" FROM assets')
    existing_numbers = [row[0] for row in c.fetchall()]
    conn.close()
    return set(existing_numbers)

def process_asset_file(uploaded_asset_file, overwrite=False):
    try:
        if uploaded_asset_file.name.endswith('.csv'):
            try:
                df = pd.read_csv(uploaded_asset_file, encoding='utf-8')
            except UnicodeDecodeError:
                df = pd.read_csv(uploaded_asset_file, encoding='cp949')
        else:
            df = pd.read_excel(uploaded_asset_file)
        
        required_columns = ['자산번호']
        if not all(col in df.columns for col in required_columns):
            st.error("파일에 '자산번호' 컬럼이 없습니다. 자산번호는 필수입니다.")
            return False
        
        existing_numbers = get_existing_asset_numbers()
        assets_to_add = df.to_dict('records')
        added_count, failed_count, duplicate_numbers = 0, 0, []
        valid_statuses = ["사용중", "미사용", "수리중", "폐기", "정상", "비정상"]
        
        for asset in assets_to_add:
            asset_data = {k: (v if pd.notna(v) else None) for k, v in asset.items()}
            if not asset_data.get('자산번호'):
                st.warning(f"자산번호가 누락된 항목 발견. 해당 행은 무시됩니다.")
                failed_count += 1
                continue
            
            status = asset_data.get('장비상태')
            if status and status not in valid_statuses:
                st.warning(f"잘못된 상태 값 '{status}' 발견 (자산번호: {asset_data['자산번호']}). '사용중'으로 대체합니다.")
                asset_data['장비상태'] = '사용중'
            
            full_asset_data = {
                '자산번호': asset_data.get('자산번호'),
                '위치': asset_data.get('위치'),
                '장비상태': asset_data.get('장비상태'),
                '제조사': asset_data.get('제조사'),
                '시리얼': asset_data.get('시리얼'),
                '구분': asset_data.get('구분'),
                '모델명': asset_data.get('모델명'),
                '기타': asset_data.get('기타', asset_data.get('시스템명'))
            }
            
            if full_asset_data['자산번호'] in existing_numbers:
                if overwrite:
                    existing_asset = next((a for a in get_assets() if a['자산번호'] == full_asset_data['자산번호']), None)
                    if existing_asset:
                        full_asset_data['id'] = existing_asset['id']
                        update_asset(full_asset_data)
                        added_count += 1
                    else:
                        duplicate_numbers.append(full_asset_data['자산번호'])
                        failed_count += 1
                else:
                    duplicate_numbers.append(full_asset_data['자산번호'])
                    failed_count += 1
                continue
            
            if add_asset(full_asset_data):
                added_count += 1
            else:
                failed_count += 1
        
        if duplicate_numbers:
            st.warning(f"다음 자산번호는 이미 DB에 존재하여 추가되지 않았습니다: {', '.join(duplicate_numbers)}")
        if added_count > 0:
            st.success(f"총 {added_count}개의 자산을 성공적으로 추가했습니다.")
        if failed_count > 0:
            st.error(f"{failed_count}개의 자산은 중복된 자산번호 또는 기타 오류로 인해 추가되지 않았습니다.")
        return added_count > 0
    except Exception as e:
        st.error(f"파일 처리 중 오류 발생: {e}")
        return False

def generate_qr_code(data):
    qr = qrcode.QRCode(version=1, box_size=10, border=5)
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill='black', back_color='white')
    buf = BytesIO()
    img.save(buf)
    buf.seek(0)
    return buf

# --- 구글 캘린더 기능 ---
def get_google_calendar_service(manual_auth=False):
    creds = None
    if os.path.exists(GOOGLE_TOKEN_PATH):
        with open(GOOGLE_TOKEN_PATH, 'rb') as token:
            creds = pickle.load(token)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            try:
                creds.refresh(Request())
            except Exception as e:
                st.warning(f"토큰 갱신 실패: {e}. 재인증이 필요합니다.")
                creds = None
        else:
            if not os.path.exists(GOOGLE_CREDS_PATH):
                st.error(f"'{GOOGLE_CREDS_PATH}' 파일을 찾을 수 없습니다. OAuth 2.0 클라이언트 ID 파일이 필요합니다.")
                return None
            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    GOOGLE_CREDS_PATH, scopes=['https://www.googleapis.com/auth/calendar.readonly']
                )
                auth_url, _ = flow.authorization_url(prompt='consent')
                st.info(f'아래 URL에 접속하여 구글 계정으로 로그인하고, 표시되는 인증 코드를 복사하여 아래에 붙여넣어 주세요.')
                st.code(auth_url)
                auth_code = st.text_input('인증 코드를 여기에 입력하세요:')
                if auth_code:
                    flow.fetch_token(code=auth_code)
                    creds = flow.credentials
                else:
                    return None
            except Exception as e:
                st.error(f"OAuth 2.0 인증 흐름 중 오류 발생: {e}")
                if os.path.exists(GOOGLE_TOKEN_PATH):
                    os.remove(GOOGLE_TOKEN_PATH)
                return None
        with open(GOOGLE_TOKEN_PATH, 'wb') as token:
            pickle.dump(creds, token)
    try:
        service = build('calendar', 'v3', credentials=creds)
        return service
    except Exception as e:
        st.error(f"Google Calendar 서비스 생성 중 오류 발생: {e}")
        return None

def get_calendar_events_from_query(query):
    try:
        today_str = datetime.now().strftime("%Y-%m-%d")
        chain_output = date_extraction_chain.invoke({"today": today_str, "question": query})
        extracted_info = json.loads(chain_output['text'])
        start_date_str = extracted_info.get("start_date", today_str)
        end_date_str = extracted_info.get("end_date", today_str)
        keyword = extracted_info.get("query", "전체")
    except Exception as e:
        return f"질문 분석 중 오류 발생: {e}. 좀 더 명확한 질문을 시도해보세요."
    service = get_google_calendar_service()
    if not service:
        return "캘린더 서비스를 사용할 수 없습니다. 먼저 인증을 완료해주세요."
    try:
        start_date_dt = datetime.fromisoformat(start_date_str)
        end_date_dt = datetime.fromisoformat(end_date_str)
        time_min = start_date_dt.isoformat() + 'Z'
        time_max = (end_date_dt + timedelta(days=1)).isoformat() + 'Z'
        calendar_list = service.calendarList().list().execute()
        all_events = []
        for calendar_entry in calendar_list.get('items', []):
            calendar_id = calendar_entry['id']
            events_result = service.events().list(
                calendarId=calendar_id, timeMin=time_min, timeMax=time_max,
                singleEvents=True, orderBy='startTime'
            ).execute()
            all_events.extend(events_result.get('items', []))
        all_events.sort(key=lambda x: x['start'].get('dateTime', x['start'].get('date')))
        events = all_events
        if keyword.lower() != '전체':
            search_terms = [term.strip().lower() for term in re.split(r'[,\s]+', keyword) if term.strip()]
            if search_terms:
                filtered_events = []
                for event in events:
                    summary = event.get('summary', '').lower()
                    if any(term in summary for term in search_terms):
                        filtered_events.append(event)
                events = filtered_events
            else:
                events = []
        if not events:
            return f"'{start_date_str}'부터 '{end_date_str}'까지 '{keyword}' 관련 일정이 없습니다."
        response_text = f"🗓️ **'{start_date_str} ~ {end_date_str}'의 '{keyword}' 관련 일정입니다.**\n\n"
        for event in events:
            start = event['start'].get('dateTime', event['start'].get('date'))
            start_dt = datetime.fromisoformat(start.replace('Z', '+00:00'))
            KST = timezone(timedelta(hours=9))
            kst_start_dt = start_dt.astimezone(KST)
            formatted_start = kst_start_dt.strftime('%Y-%m-%d %H:%M') if 'dateTime' in event['start'] else kst_start_dt.strftime('%Y-%m-%d') + " (종일)"
            summary = event['summary']
            response_text += f"- **{summary}** ({formatted_start})\n"
        return response_text
    except Exception as e:
        return f"일정 조회 중 오류 발생: {e}"

# --- RAG 기능 ---
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def parse_excel(file_path):
    try:
        if file_path.endswith('.csv'):
            try:
                df = pd.read_csv(file_path, encoding='utf-8')
            except UnicodeDecodeError:
                df = pd.read_csv(file_path, encoding='cp949')
        else:
            df = pd.read_excel(file_path)
    except Exception as e:
        st.error(f"파일 파싱 오류: {e}")
        return []
    chunks = []
    for _, row in df.iterrows():
        text = " | ".join([str(val) for val in row.values if pd.notna(val)])
        metadata = {"file": os.path.basename(file_path), "type": "excel"}
        for key, value in row.to_dict().items():
            if pd.notna(value):
                metadata[str(key)] = str(value)
        chunks.append({"text": text, "metadata": metadata})
    st.write(f"Parsed {len(chunks)} chunks from {file_path}")
    return chunks

def parse_docx(file_path):
    try:
        doc = Document(file_path)
    except Exception as e:
        st.error(f"DOCX 파일 파싱 오류: {e}")
        return []
    chunks = []
    for i, table in enumerate(doc.tables):
        table_text = ""
        for row in table.rows:
            row_text = " | ".join([cell.text.strip() for cell in row.cells])
            table_text += row_text + "\n"
        if table_text.strip():
            chunks.append({
                "text": table_text.strip(),
                "metadata": {"file": os.path.basename(file_path), "type": "docx_table", "table_index": i}
            })
    max_chunk_size = 500
    current_section = "본문"
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if text.startswith(("1.", "2.", "3.", "4.", "5.", "6.", "7.", "8.")):
            current_section = text.split()[0]
        while len(text) > max_chunk_size:
            split_point = text.rfind('.', 0, max_chunk_size)
            if split_point == -1:
                split_point = max_chunk_size
            chunk_text = text[:split_point]
            chunks.append({
                "text": chunk_text,
                "metadata": {"file": os.path.basename(file_path), "type": "docx_paragraph", "section": current_section}
            })
            text = text[split_point:].strip()
        if text:
            chunks.append({
                "text": text,
                "metadata": {"file": os.path.basename(file_path), "type": "docx_paragraph", "section": current_section}
            })
    st.write(f"Parsed {len(chunks)} chunks from {file_path}")
    return chunks

def parse_pdf(file_path):
    try:
        with pdfplumber.open(file_path) as pdf:
            chunks = []
            max_chunk_size = 500
            for page in pdf.pages:
                text = page.extract_text()
                tables = page.extract_tables()
                for table in tables:
                    table_text = "\n".join([" | ".join(str(cell) if cell else "" for cell in row) for row in table])
                    if table_text.strip():
                        chunks.append({
                            "text": table_text,
                            "metadata": {"file": os.path.basename(file_path), "type": "pdf_table", "page": page.page_number}
                        })
                if text:
                    while len(text) > max_chunk_size:
                        split_point = text.rfind('。', 0, max_chunk_size)
                        if split_point == -1: split_point = text.rfind('.', 0, max_chunk_size)
                        if split_point == -1: split_point = max_chunk_size
                        chunk_text = text[:split_point + 1]
                        chunks.append({
                            "text": chunk_text.strip(),
                            "metadata": {"file": os.path.basename(file_path), "type": "pdf_text", "page": page.page_number}
                        })
                        text = text[split_point + 1:].strip()
                    if text:
                        chunks.append({
                            "text": text,
                            "metadata": {"file": os.path.basename(file_path), "type": "pdf_text", "page": page.page_number}
                        })
    except Exception as e:
        st.error(f"PDF 파일 파싱 오류: {e}")
        return []
    st.write(f"Parsed {len(chunks)} chunks from {file_path}")
    return chunks

from pptx import Presentation
import streamlit as st
import os

def parse_pptx(file_path):
    try:
        prs = Presentation(file_path)
        chunks = []
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            if slide_text.strip():
                chunks.append({
                    "text": slide_text.strip(),
                    "metadata": {"file": os.path.basename(file_path), "type": "pptx_slide", "slide_number": prs.slides.index(slide) + 1}
                })
    except Exception as e:
        st.error(f"PPTX 파일 파싱 오류: {e}")
        return []
    st.write(f"Parsed {len(chunks)} chunks from {file_path}")
    return chunks

def parse_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        chunks = []
        max_chunk_size = 500
        while len(text) > max_chunk_size:
            split_point = text.rfind('。', 0, max_chunk_size)
            if split_point == -1: split_point = text.rfind('.', 0, max_chunk_size)
            if split_point == -1: split_point = max_chunk_size
            chunk_text = text[:split_point + 1]
            chunks.append({"text": chunk_text.strip(), "metadata": {"file": os.path.basename(file_path), "type": "txt"}})
            text = text[split_point + 1:].strip()
        if text:
            chunks.append({"text": text, "metadata": {"file": os.path.basename(file_path), "type": "txt"}})
    except Exception as e:
        st.error(f"TXT 파일 파싱 오류: {e}")
        return []
    st.write(f"Parsed {len(chunks)} chunks from {file_path}")
    return chunks

def parse_image(file_path):

    try:
        img = Image.open(file_path)
        img = img.convert('L')
        img = ImageEnhance.Contrast(img).enhance(2.0)
        img = img.filter(ImageFilter.SHARPEN)
        text = pytesseract.image_to_string(img, lang='kor+eng', config='--psm 6')
        if not text.strip():
            st.warning(f"{file_path}에서 텍스트를 추출하지 못했습니다.")
            return []
        chunks = []
        max_chunk_size = 500
        while len(text) > max_chunk_size:
            split_point = text.rfind('。', 0, max_chunk_size)
            if split_point == -1: split_point = text.rfind('.', 0, max_chunk_size)
            if split_point == -1: split_point = max_chunk_size
            chunk_text = text[:split_point + 1]
            chunks.append({"text": chunk_text.strip(), "metadata": {"file": os.path.basename(file_path), "type": "image"}})
            text = text[split_point + 1:].strip()
        if text:
            chunks.append({"text": text, "metadata": {"file": os.path.basename(file_path), "type": "image"}})
    except Exception as e:
        st.error(f"이미지 파일 파싱 오류: {e}")
        return []
    st.write(f"Parsed {len(chunks)} chunks from {file_path}")
    return chunks

def embed_and_store(chunks):
    if not chunks:
        return
    co_client, faiss_index, _, current_metadata = get_clients()
    docs_to_add = [chunk["text"] for chunk in chunks]
    metadatas_to_add = [chunk for chunk in chunks]
    try:
        embeddings = co_client.embed(texts=docs_to_add, model="embed-multilingual-v3.0", input_type="search_document").embeddings
        embeddings = np.array(embeddings).astype('float32')
        faiss_index.add(embeddings)
        current_metadata.extend(metadatas_to_add)
        faiss.write_index(faiss_index, FAISS_INDEX_PATH)
        with open(FAISS_METADATA_PATH, 'wb') as f:
            pickle.dump(current_metadata, f)
        st.success(f"Stored {len(chunks)} chunks in FAISS")
    except Exception as e:
        st.error(f"임베딩 저장 오류: {e}")

def hybrid_retriever(query, top_k=5):
    co_client, faiss_index, _, current_metadata = get_clients()
    try:
        query_embedding = co_client.embed(texts=[query], model="embed-multilingual-v3.0", input_type="search_query").embeddings[0]
        query_embedding = np.array([query_embedding]).astype('float32')
        distances, indices = faiss_index.search(query_embedding, top_k * 2)
        retrieved_docs = []
        retrieved_metadatas = []
        for idx in indices[0]:
            if idx < len(current_metadata) and 'text' in current_metadata[idx]:
                retrieved_docs.append(current_metadata[idx]["text"])
                retrieved_metadatas.append(current_metadata[idx]["metadata"])
        if not retrieved_docs:
            st.warning("No documents found for reranking.")
            return [], []
        rerank_results = co_client.rerank(
            query=query, 
            documents=retrieved_docs, 
            top_n=top_k, 
            model="rerank-multilingual-v3.0",
            return_documents=True
        )
        reranked_docs = []
        reranked_metadatas = []
        if rerank_results.results:
            for result in rerank_results.results:
                if result.document:
                    reranked_docs.append(result.document.text)
                    reranked_metadatas.append(retrieved_metadatas[result.index])
        st.write(f"Retrieved {len(reranked_docs)} documents after reranking")
        return reranked_docs, reranked_metadatas
    except Exception as e:
        st.error(f"리트리버 오류: {e}")
        return [], []

def generate_narrative_response(query, docs, metadatas):
    context = "\n\n---\n\n".join(docs)
    prompt = f"""
    당신은 제공된 '문서'의 내용을 바탕으로 사용자의 질문에 대해 친절하고 상세하게 서술형으로 답변하는 AI 어시스턴트입니다.

    # 지시사항
    1. **질문 파악:** 사용자의 질문('{query}')의 핵심 의도를 정확히 파악합니다.
    2. **문서 분석:** 제공된 '문서' 내용을 깊이 있게 분석하여 질문과 관련된 정보를 모두 찾습니다.
    3. **답변 생성:**
        - 분석한 정보를 바탕으로, 자연스러운 한국어 문장으로 답변을 구성합니다.
        - 딱딱한 형식이 아닌, 사람이 설명해주듯이 친절한 어조를 사용해주세요.
        - 필요하다면, 정보를 목록이나 단락으로 나누어 가독성을 높이세요.
    4. **추측 금지:** 문서에 없는 내용은 절대로 답변에 포함하지 마세요. 정보가 부족할 경우, "제공된 문서의 내용만으로는 질문에 답변하기 어렵습니다." 라고 솔직하게 답변해주세요.

    # 문서 (이 내용을 기반으로 답변을 생성하세요)
    ---
    {context}
    ---

    # 질문
    {query}

    # 답변
    """
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that provides detailed, narrative answers based on provided documents."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"GPT-4o 호출 오류: {e}")
        return "답변 생성 중 오류가 발생했습니다."

def classify_query_intent(query):
    prompt = f"""
    사용자의 질문을 분석하여, 답변이 '테이블(표)' 형식으로 제공되어야 하는지, 아니면 '서술형'으로 제공되어야 하는지 분류해주세요.

    # 지시사항
    1. **'table'로 분류해야 하는 경우 (매우 제한적):**
        - 명시적으로 '목록', '리스트', '현황', '장비', '표', '테이블' 등 표 형식의 데이터를 요청하는 경우.
        - 여러 항목에 대한 비교나 요약을 요청하며, 그 결과가 명확히 구조화된 데이터 형태일 때.
        - 예시: "A 장비 목록 보여줘", "HRS코리아 유지보수 현황 알려줘", "모든 자산 리스트", "2024년 네트워크 장비 목록", "이 데이터를 표로 보여줘"

    2. **'narrative'로 분류해야 하는 경우 (기본값):**
        - 설명, 방법, 이유, 정의 등 서술적인 답변이 필요한 모든 질문.
        - 단일 정보에 대한 질문.
        - 일반적인 대화형 질문.
        - 'table'로 분류될 명확한 근거가 없는 모든 질문은 'narrative'로 분류합니다.
        - 예시: "PAS-K가 뭐야?", "와이파이 접속 어떻게 해?", "어제 무슨 일 있었어?", "백종윤의 일정이 어떻게 돼?", "이 문서의 요약은?", "이 개념에 대해 설명해줘"

    3. **출력 형식:**
        - 분류 결과를 오직 'table' 또는 'narrative' 단어 하나로만 반환해야 합니다. 다른 설명은 절대 추가하지 마세요

    # 질문
    {query}

    # 분류 결과
    """
    try:
        response = openai_client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a query classifier. Your only output should be 'table' or 'narrative'."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.0,
            max_tokens=5
        )
        classification = response.choices[0].message.content.strip().lower()
        if classification not in ['table', 'narrative']:
            return 'narrative'
        return classification
    except Exception as e:
        st.error(f"Query classification error: {e}")
        return 'narrative'

def generate_response(query, docs, metadatas):
    context = "\n\n---\n\n".join(docs)
    prompt = f"""
    당신은 제공된 '문서'의 내용을 분석하여, 사용자의 질문에 대한 답변을 JSON 형식으로 추출하는 AI입니다.

    # 지시사항
    1. **질문 분석:** 사용자의 질문('{query}')을 정확히 이해합니다.
    2. **데이터 추출:** '문서' 내용 전체를 꼼꼼히 분석하여, 질문에 해당하는 모든 데이터를 찾습니다. 문서가 CSV 표 형식인 경우, 모든 행(row)을 빠짐없이 확인해야 합니다.
    3. **JSON 형식화:**
        - 추출한 데이터를 JSON 형식으로 만듭니다.
        - JSON의 최상위 키는 'data' 여야 하고, 그 값은 객체들의 배열(list of objects)이어야 합니다.
        - 각 객체는 문서의 한 행(row)에 해당하며, 문서의 헤더(header)를 키(key)로 사용해야 합니다.
        - 예시: {{"data": [{{"column1": "value1", "column2": "value2"}}, {{"column1": "value3", "column2": "value4"}}]}}
    4. **완전성:** 질문이 '전체 목록'이나 '모든 품목'을 요구하는 경우, 단 하나의 데이터도 누락해서는 안 됩니다. 문서에 있는 모든 관련 데이터를 JSON에 포함시켜야 합니다.
    5. **추측 금지:** 문서에 없는 내용은 절대로 추가하지 마세요. 정보가 없으면 {{"data": []}} 와 같이 빈 배열을 반환합니다.

    # 문서 (이 내용을 기반으로 JSON을 생성하세요)
    ---
    {context}
    ---
    """
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a precise AI that extracts information from documents and returns it in JSON format."},
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"},
            temperature=0.0,
            max_tokens=16000
        )
        st.session_state.json_response = response.choices[0].message.content
        return st.session_state.json_response
    except Exception as e:
        st.error(f"GPT-4o 호출 오류: {e}")
        return None

# --- Streamlit UI ---
st.title("서비스사업부 챗봇 시스템")
init_asset_db()
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["업무조회", "업무일정", "자산현황", "업무학습", "백터DB", "장비관리"])

# --- 업무조회 탭 ---
with tab1:
    st.header("AI 업무 챗봇")
    with st.form(key="query_form"):
        query = st.text_input("유지보수 업무와 관련하여 물어보세요 (예: 한국국제협력단 유지보수 현황):", key="query")
        submit_button = st.form_submit_button("질문하기")
    if submit_button and query:
        with st.spinner("질문 의도를 파악하고 답변을 생성하는 중입니다..."):
            docs, metadatas = hybrid_retriever(query)
            if not docs:
                st.warning("관련 문서를 찾을 수 없습니다.")
            else:
                intent = classify_query_intent(query)
                if intent == 'table':
                    st.info("사용자의 질문 의도를 '테이블'로 파악했습니다. 표 형식으로 답변을 생성합니다.")
                    json_response_str = generate_response(query, docs, metadatas)
                    if json_response_str:
                        try:
                            response_data = json.loads(json_response_str)
                            data = response_data.get('data', [])
                            if data and isinstance(data, list):
                                st.subheader("검색 결과")
                                df = pd.DataFrame(data)
                                st.dataframe(df, use_container_width=True, hide_index=True)
                            else:
                                st.warning("표로 만들 데이터를 찾지 못했습니다. 대신 서술형 답변을 생성합니다.")
                                narrative_response = generate_narrative_response(query, docs, metadatas)
                                st.subheader("AI 응답")
                                st.markdown(narrative_response)
                        except json.JSONDecodeError:
                            st.error("JSON 응답을 파싱하는 데 실패했습니다. 원본 응답을 표시합니다.")
                            st.write(json_response_str)
                    else:
                        st.warning("AI로부터 테이블 형식의 응답을 받지 못했습니다.")
                else:
                    st.info("사용자의 질문 의도를 '서술형'으로 파악했습니다. 대화 형태로 답변을 생성합니다.")
                    narrative_response = generate_narrative_response(query, docs, metadatas)
                    st.subheader("AI 응답")
                    st.markdown(narrative_response)
    elif submit_button and not query:
        st.warning("질문을 입력하세요.")

# --- 업무일정 탭 ---
with tab2:
    st.header("📅 업무일정 조회")
    if not os.path.exists(GOOGLE_TOKEN_PATH):
        st.info("업무일정 조회를 위해 구글 캘린더 연동이 필요합니다.")
        if st.button("구글 캘린더 연동 및 인증"):
            with st.spinner("인증을 진행합니다. 터미널의 안내에 따라 브라우저에서 인증을 완료해주세요..."):
                get_google_calendar_service(manual_auth=True)
                st.success("인증이 완료되었습니다! 이제 일정을 조회할 수 있습니다.")
                st.rerun()
    else:
        with st.form(key="schedule_form"):
            schedule_query = st.text_input("업무 일정을 자연어로 물어보세요:", placeholder="예: 8월 1일 ~ 8월 10일 현우, 원준 일정")
            schedule_submit_button = st.form_submit_button("일정 검색")
        if schedule_submit_button and schedule_query:
            with st.spinner("일정을 분석하고 있습니다..."):
                events_response = get_calendar_events_from_query(schedule_query)
                st.markdown(events_response)

# --- 자산현황 탭 ---
with tab3:
    st.header("📋 자산 현황 및 관리")
    st.subheader("자산 목록")
    assets = get_assets()
    if assets:
        df = pd.DataFrame(assets)
        st.dataframe(df, use_container_width=True, hide_index=True)
    else:
        st.info("등록된 자산이 없습니다.")
    
    st.markdown("---")
    with st.expander("➕ 새 자산 추가하기"):
        with st.form(key="asset_add_form", clear_on_submit=True):
            c1, c2 = st.columns(2)
            asset_data = {
                "자산번호": c1.text_input("자산번호*"),
                "위치": c2.text_input("위치"),
                "장비상태": c1.selectbox("장비상태", ["사용중", "미사용", "수리중", "폐기", "정상", "비정상"]),
                "제조사": c2.text_input("제조사"),
                "시리얼": c1.text_input("시리얼"),
                "구분": c2.text_input("구분 (예: 서버, 노트북)"),
                "모델명": c1.text_input("모델명"),
                "기타": c2.text_input("기타")
            }
            if st.form_submit_button("자산 추가"):
                if not asset_data["자산번호"]:
                    st.error("자산번호는 필수 항목입니다.")
                else:
                    if add_asset(asset_data):
                        st.success("자산이 성공적으로 추가되었습니다.")
                        st.rerun()

    with st.expander("✏️ 자산 수정 또는 삭제하기"):
        if assets:
            selected_id = st.selectbox("관리할 자산의 ID를 선택하세요.", options=[a['id'] for a in assets], format_func=lambda x: f"ID: {x}")
            selected_asset = next((a for a in assets if a['id'] == selected_id), None)
            if selected_asset:
                col1, col2 = st.columns([3, 1])
                with col1:
                    with st.form(key=f"edit_form_{selected_id}"):
                        st.write(f"**ID {selected_id} 자산 정보 수정**")
                        status_options = get_unique_statuses()
                        current_status = selected_asset['장비상태'] if selected_asset['장비상태'] else "사용중"
                        if current_status not in status_options:
                            status_options.append(current_status)
                        try:
                            status_index = status_options.index(current_status)
                        except ValueError:
                            status_index = 0
                            st.warning(f"상태 '{current_status}'가 유효하지 않습니다. 기본값 '사용중'으로 설정됩니다.")
                        updated_data = {
                            "id": selected_id,
                            "자산번호": st.text_input("자산번호", value=selected_asset['자산번호']),
                            "위치": st.text_input("위치", value=selected_asset['위치']),
                            "장비상태": st.selectbox("장비상태", status_options, index=status_index),
                            "제조사": st.text_input("제조사", value=selected_asset['제조사']),
                            "시리얼": st.text_input("시리얼", value=selected_asset['시리얼']),
                            "구분": st.text_input("구분", value=selected_asset['구분']),
                            "모델명": st.text_input("모델명", value=selected_asset['모델명']),
                            "기타": st.text_input("기타", value=selected_asset['기타'])
                        }
                        if st.form_submit_button("수정 완료"):
                            update_asset(updated_data)
                            st.success(f"ID {selected_id} 자산 정보가 수정되었습니다.")
                            st.rerun()
                with col2:
                    st.write("**QR 코드 & 삭제**")
                    qr_data_str = ", ".join([f"{k}: {v}" for k, v in selected_asset.items() if k != 'id'])
                    qr_img = generate_qr_code(qr_data_str)
                    st.image(qr_img, caption="자산 정보 QR 코드")
                    if st.button("이 자산 삭제", key=f"delete_{selected_id}", type="primary"):
                        delete_asset(selected_id)
                        st.warning(f"ID {selected_id} 자산이 삭제되었습니다.")
                        st.rerun()
        else:
            st.info("수정 또는 삭제할 자산이 없습니다.")
            
    st.markdown("---")
    st.subheader("🗂️ 자산 일괄 등록 및 DB 관리")
    if 'asset_upload_processed' not in st.session_state:
        st.session_state.asset_upload_processed = False
    overwrite_option = st.checkbox("중복 자산번호가 있으면 덮어쓰기")
    uploaded_asset_file = st.file_uploader("자산 목록 엑셀/CSV 파일 업로드", type=["xlsx", "csv"], key="asset_uploader")
    if uploaded_asset_file and not st.session_state.asset_upload_processed:
        with st.spinner("파일을 처리하여 자산을 등록하는 중..."):
            success = process_asset_file(uploaded_asset_file, overwrite=overwrite_option)
            if success:
                st.session_state.asset_upload_processed = True
                st.rerun()
            else:
                st.session_state.asset_upload_processed = False
    if st.session_state.asset_upload_processed:
        st.session_state.asset_upload_processed = False

    if 'confirm_asset_reset' not in st.session_state:
        st.session_state.confirm_asset_reset = False
    if st.button("🗑️ 모든 자산 DB 초기화", type="primary"):
        st.session_state.confirm_asset_reset = True
    if st.session_state.confirm_asset_reset:
        st.warning("모든 자산 데이터가 영구적으로 삭제됩니다. 계속하시겠습니까?", icon="⚠️")
        c1, c2 = st.columns(2)
        if c1.button("예, 초기화합니다.", key="confirm_asset_delete"):
            try:
                delete_all_assets()
                st.session_state.confirm_asset_reset = False
                st.success("모든 자산 데이터가 성공적으로 초기화되었습니다.")
                st.rerun()
            except Exception as e:
                st.error(f"초기화 중 오류 발생: {e}")
                st.session_state.confirm_asset_reset = False
        if c2.button("아니요, 취소합니다.", key="cancel_asset_delete"):
            st.session_state.confirm_asset_reset = False
            st.rerun()

# --- 업무학습 탭 ---
with tab4:
    st.header("업무학습")
    uploaded_file = st.file_uploader("문서를 업로드하세요(예: CSV, DOCX, PDF, JPG, PNG, PPT, TXT..)", type=["csv", "docx", "pdf", "jpg", "jpeg", "png", "ppt", "pptx", "txt"])
    if uploaded_file is not None:
        if uploaded_file.size > MAX_FILE_SIZE:
            st.error("파일 크기는 16MB를 초과할 수 없습니다.")
        elif allowed_file(uploaded_file.name):
            file_path = os.path.join(UPLOAD_FOLDER, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            try:
                if uploaded_file.name.endswith('.csv'): chunks = parse_excel(file_path)
                elif uploaded_file.name.endswith('.docx'): chunks = parse_docx(file_path)
                elif uploaded_file.name.endswith('.pdf'): chunks = parse_pdf(file_path)
                elif uploaded_file.name.endswith(('.jpg', '.jpeg', 'png')): chunks = parse_image(file_path)
                elif uploaded_file.name.endswith(('.ppt', '.pptx')): chunks = parse_pptx(file_path)
                elif uploaded_file.name.endswith('.txt'): chunks = parse_txt(file_path)
                else:
                    st.error("지원하지 않는 파일 형식입니다.")
                    chunks = []
                if chunks:
                    embed_and_store(chunks)
                    st.success(f"{uploaded_file.name} 업로드 및 처리 완료!")
                else:
                    st.error("파일에서 데이터를 추출하지 못했습니다.")
            except Exception as e:
                st.error(f"파일 처리 중 오류: {e}")
        else:
            st.error("CSV, DOCX, PDF, JPG, PNG 파일만 업로드 가능합니다.")

# --- 백터DB 탭 ---
with tab5:
    st.header("🗄️ 자산관리 백터DB")
    if 'confirm_vector_db_reset' not in st.session_state:
        st.session_state.confirm_vector_db_reset = False
    if st.button("자산관리 DB 초기화", help="학습된 모든 문서 데이터를 삭제합니다.", type="primary"):
        st.session_state.confirm_vector_db_reset = True
    if st.session_state.confirm_vector_db_reset:
        st.warning("모든 FAISS 벡터 데이터가 영구적으로 삭제됩니다. 계속하시겠습니까?", icon="⚠️")
        c1, c2 = st.columns(2)
        if c1.button("예, 초기화합니다.", key="confirm_vector_db_delete"):
            try:
                dimension = 1024
                new_index = faiss.IndexFlatL2(dimension)
                new_metadata = []
                faiss.write_index(new_index, FAISS_INDEX_PATH)
                with open(FAISS_METADATA_PATH, 'wb') as f:
                    pickle.dump(new_metadata, f)
                get_clients.clear()
                st.success("FAISS 인덱스 및 메타데이터 초기화 완료!")
                st.session_state.confirm_vector_db_reset = False
                st.rerun()
            except Exception as e:
                st.error(f"FAISS 초기화 오류: {e}")
                st.session_state.confirm_vector_db_reset = False
        if c2.button("아니요, 취소합니다.", key="cancel_vector_db_delete"):
            st.session_state.confirm_vector_db_reset = False
            st.rerun()
    _, _, _, current_metadata = get_clients()
    if current_metadata:
        st.subheader("저장된 문서 Chunk 목록")
        doc_count = len(current_metadata)
        st.write(f"총 {doc_count}개의 문서 Chunk가 저장되어 있습니다. (최대 10개 표시)")
        for i, meta in enumerate(current_metadata):
            if i >= 10:
                break
            if not isinstance(meta, dict) or 'text' not in meta:
                st.warning(f"Chunk {i+1}에 잘못된 메타데이터 형식이 있습니다: {meta}")
                continue
            with st.expander(f"**Chunk {i+1}**: {meta['text'][:80]}..."):
                st.write(f"**전체 내용**: {meta['text']}")
                st.write(f"**메타데이터**: {meta['metadata']}")
    else:
        st.info("저장된 학습 데이터가 없습니다")

# --- 장비관리 탭 ---
with tab6:
    st.header("📡 네트워크 장비 관리")
    subtab1, subtab2 = st.tabs(["데이터 수집 및 설정", "네트워크 챗봇"])
    
    with subtab1:
        st.subheader("⚙️ 데이터 수집 설정")
        with st.expander("1. 장비 접속 정보", expanded=True):
            network_host = st.text_input("Host IP", value="192.168.9.19", key="network_host")
            network_username = st.text_input("Username", value="inner", key="network_username")
            network_password = st.text_input("Password", type="password", value="", key="network_password")
            network_secret = st.text_input("Enable Secret", type="password", value="", key="network_secret")
            network_port = st.number_input("SSH Port", min_value=1, max_value=65535, value=22, key="network_port")
            ssh_key_file = st.text_input("SSH Key File Path (Optional)", placeholder="/path/to/your/private_key", key="ssh_key_file")
        
        with st.expander("2. 명령어 목록 설정"):
            uploaded_file = st.file_uploader("commands.txt 파일 업로드", type="txt", key="network_commands_uploader")
            default_commands = [
                "show version", "show running-config", "show log", "show ip interface brief",
                "show vlan", "show interface status", "show arp", "show mac address-table"
            ]
            if uploaded_file is not None:
                try:
                    commands_to_run = uploaded_file.getvalue().decode("utf-8").splitlines()
                    st.session_state.network_commands = [cmd for cmd in commands_to_run if cmd.strip()]
                except Exception as e:
                    st.error(f"파일 읽기 오류: {e}")
                    st.session_state.network_commands = default_commands
            else:
                st.session_state.network_commands = st.session_state.get('network_commands', default_commands)
            st.text_area("적용될 명령어 목록", value='\n'.join(st.session_state.network_commands), height=200, disabled=True)
        
        if st.button("🔄 데이터 수집 및 인덱스 생성", type="primary"):
            # 데이터 수집 전 항상 기존 인덱스 파일 삭제
            try:
                if os.path.exists(NETWORK_FAISS_INDEX_PATH):
                    import shutil
                    shutil.rmtree(NETWORK_FAISS_INDEX_PATH)
                if os.path.exists(NETWORK_RAW_DATA_FILE):
                    os.remove(NETWORK_RAW_DATA_FILE)
                load_network_rag_chain.clear()
                st.toast("기존 네트워크 벡터 DB를 초기화했습니다.")
            except Exception as e:
                st.error(f"기존 네트워크 DB 삭제 중 오류: {e}")
                st.stop()

            device_info = {
                "device_type": "cisco_ios",
                "host": st.session_state.network_host,
                "username": st.session_state.network_username,
                "password": st.session_state.network_password,
                "secret": st.session_state.network_secret,
                "port": st.session_state.network_port,
            }
            if st.session_state.ssh_key_file:
                device_info['use_keys'] = True
                device_info['key_file'] = st.session_state.ssh_key_file
            with st.spinner("장비에서 데이터를 수집하고 인덱스를 다시 생성하는 중..."):
                success, message = fetch_network_data(device_info, st.session_state.network_commands)
                if success:
                    success, message = build_network_vector_store()
            if success:
                st.success("데이터 수집 및 인덱스 업데이트 완료!")
                if 'network_messages' in st.session_state:
                    del st.session_state['network_messages']
                load_network_rag_chain.clear()
                st.rerun()
            else:
                st.error(message)
        
        if 'confirm_network_db_reset' not in st.session_state:
            st.session_state.confirm_network_db_reset = False
        if st.button("네트워크 백터DB 초기화", help="네트워크 관련 학습 데이터를 삭제합니다.", type="primary"):
            st.session_state.confirm_network_db_reset = True
        if st.session_state.confirm_network_db_reset:
            st.warning("모든 네트워크 FAISS 벡터 데이터가 영구적으로 삭제됩니다. 계속하시겠습니까?", icon="⚠️")
            c1, c2 = st.columns(2)
            if c1.button("예, 초기화합니다.", key="confirm_network_db_delete"):
                try:
                    if os.path.exists(NETWORK_FAISS_INDEX_PATH):
                        import shutil
                        shutil.rmtree(NETWORK_FAISS_INDEX_PATH)
                    if os.path.exists(NETWORK_RAW_DATA_FILE):
                        os.remove(NETWORK_RAW_DATA_FILE)
                    load_network_rag_chain.clear()
                    st.success("네트워크 FAISS 인덱스 및 데이터 초기화 완료!")
                except Exception as e:
                    st.error(f"네트워크 FAISS 초기화 오류: {e}")
                finally:
                    st.session_state.confirm_network_db_reset = False
                    st.rerun()
            if c2.button("아니요, 취소합니다.", key="cancel_network_db_delete"):
                st.session_state.confirm_network_db_reset = False
                st.rerun()
    
    with subtab2:
        st.subheader("🤖 네트워크 챗봇")
        if not os.path.exists(NETWORK_FAISS_INDEX_PATH):
            st.info("아직 네트워크 데이터가 수집되지 않았습니다. '데이터 수집 및 설정' 탭에서 데이터를 수집하세요.")
        else:
            qa_chain = load_network_rag_chain()
            if not qa_chain:
                st.error("챗봇을 로드할 수 없습니다. 데이터를 다시 수집해주세요.")
            else:
                if "network_messages" not in st.session_state:
                    st.session_state.network_messages = []
                for message in st.session_state.network_messages:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"])
                if prompt := st.chat_input("예: VLAN 9번에 할당된 포트는 무엇인가요?"):
                    st.session_state.network_messages.append({"role": "user", "content": prompt})
                    with st.chat_message("user"):
                        st.markdown(prompt)
                    with st.chat_message("assistant"):
                        with st.spinner("답변을 생성하는 중..."):
                            response = qa_chain.invoke({"query": prompt})
                            answer = response['result']
                            st.markdown(answer)
                            st.session_state.network_messages.append({"role": "assistant", "content": answer})
        
        st.markdown("---")
        st.subheader("📄 수집된 네트워크 Raw Data")
        try:
            with open(NETWORK_RAW_DATA_FILE, "r", encoding="utf-8") as f:
                raw_data = f.read()
            st.code(raw_data, language="log")
        except FileNotFoundError:
            st.warning("아직 수집된 데이터가 없습니다. 먼저 데이터 수집을 실행해주세요.")
        except Exception as e:
            st.error(f"로우 데이터를 읽는 중 오류가 발생했습니다: {e}")
    
    